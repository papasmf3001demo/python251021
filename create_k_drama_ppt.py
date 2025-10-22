#pip install python-pptx Pillow requests
import os
import requests
from urllib.parse import quote
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image
from io import BytesIO

OUT_DIR = r"C:\work"
OUT_FILE = os.path.join(OUT_DIR, "k-drama.pptx")  # .ppt 대신 .pptx 권장
TMP_IMG_DIR = os.path.join(OUT_DIR, "tmp_images")
os.makedirs(TMP_IMG_DIR, exist_ok=True)

DRAMA_TITLES = [
    "Crash Landing on You",
    "Squid Game",
    "Descendants of the Sun",
    "Goblin (TV series)",
    "My Love from the Star"
]

def fetch_wikipedia_summary(title):
    url = "https://en.wikipedia.org/api/rest_v1/page/summary/" + quote(title)
    r = requests.get(url, timeout=10)
    if r.status_code != 200:
        return None
    data = r.json()
    summary = data.get("extract")
    image_url = None
    # REST summary may include originalimage
    if "originalimage" in data:
        image_url = data["originalimage"].get("source")
    # fallback: thumbnail
    elif "thumbnail" in data:
        image_url = data["thumbnail"].get("source")
    return {"title": data.get("title", title), "summary": summary, "image": image_url}

def download_image(url, name):
    if not url:
        return None
    try:
        r = requests.get(url, timeout=10)
        r.raise_for_status()
        img = Image.open(BytesIO(r.content)).convert("RGB")
        path = os.path.join(TMP_IMG_DIR, name)
        img.save(path, format="JPEG")
        return path
    except Exception:
        return None

def create_ppt(drama_infos, out_path):
    prs = Presentation()
    # 제목 슬라이드
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "K-Drama Highlights"
    if slide.placeholders:
        try:
            slide.placeholders[1].text = "자동 생성 자료 — 출처: Wikipedia"
        except Exception:
            pass

    for info in drama_infos:
        layout = prs.slide_layouts[5]  # blank
        s = prs.slides.add_slide(layout)
        # 제목 상자
        left = Inches(0.5)
        top = Inches(0.3)
        width = Inches(6)
        height = Inches(1)
        title_box = s.shapes.add_textbox(left, top, width, height)
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = info.get("title", "")
        p.font.size = Pt(28)

        # 요약 상자
        left = Inches(0.5)
        top = Inches(1.2)
        width = Inches(5.5)
        height = Inches(3.5)
        summary_box = s.shapes.add_textbox(left, top, width, height)
        tf2 = summary_box.text_frame
        summary = info.get("summary") or "요약을 찾을 수 없습니다."
        tf2.text = summary
        for paragraph in tf2.paragraphs:
            paragraph.font.size = Pt(12)

        # 이미지 넣기 (오른쪽)
        img_path = info.get("image_path")
        if img_path and os.path.exists(img_path):
            img_left = Inches(6.2)
            img_top = Inches(1.2)
            # 이미지 크기 조정: 최대 너비/높이 설정
            s.shapes.add_picture(img_path, img_left, img_top, width=Inches(3))

    prs.save(out_path)
    print("Saved:", out_path)

def main():
    infos = []
    for t in DRAMA_TITLES:
        data = fetch_wikipedia_summary(t)
        if not data:
            data = {"title": t, "summary": None, "image": None}
        img_path = None
        if data.get("image"):
            safe_name = quote(data["title"].replace(" ", "_")) + ".jpg"
            img_path = download_image(data["image"], safe_name)
        data["image_path"] = img_path
        infos.append(data)

    create_ppt(infos, OUT_FILE)
    # 임시 이미지 정리(원하면 주석 처리)
    # for f in os.listdir(TMP_IMG_DIR):
    #     os.remove(os.path.join(TMP_IMG_DIR, f))

if __name__ == "__main__":
    main()